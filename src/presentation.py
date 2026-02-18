import asyncio

from pptx import Presentation
from pptx.util import Pt

from src.api import PresentationAPI


class UserPresentation:
    def __init__(self, theme, slides_count):
        self.theme = theme
        self.slides_count = slides_count
        self.api = PresentationAPI()
        self.prs = Presentation()

    async def create_presentation(self, output_filename="presentation.pptx"):
        # 1. Генерируем план
        topics = await self.api.generate_topics(self.theme, self.slides_count)

        # 2. Титульный слайд
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = self.theme.upper()
        slide.placeholders[1].text = "ai generated content :))"

        # 3. Основные слайды
        for num, topic in topics.items():
            print(f"Слайд {num}: {topic}")

            # Получаем контент в виде списка строк
            points = await self.api.generate_slide_content(self.theme, topic)

            # Создаем слайд
            slide_layout = self.prs.slide_layouts[1]  # Заголовок и текст
            slide = self.prs.slides.add_slide(slide_layout)

            # Заголовок
            slide.shapes.title.text = topic

            # Текстовый блок
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = ""  # Очищаем дефолтный текст

            for point in points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0  # Уровень списка (0 = обычный маркер)
                p.font.size = Pt(20)
                p.space_after = Pt(12)

        # 4. Сохранение
        self.prs.save(output_filename)
        print(f"\nФайл '{output_filename}' успешно создан!")


async def main():
    theme = "Влияние экологии на здоровье человека"
    count = 4
    pres = UserPresentation(theme, count)
    await pres.create_presentation("Ecology_Presentation.pptx")


if __name__ == "__main__":
    asyncio.run(main())
